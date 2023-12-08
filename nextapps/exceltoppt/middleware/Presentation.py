from functions import api_openai
from dotenv import load_dotenv

from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
from difflib import get_close_matches
from utils import load_json
from utils import extract_metrics_and_dimensions


import re
import os
import openai
import openpyxl
import asyncio


load_dotenv()
openai.api_key = os.getenv("OPENAI_API_KEY")
os.environ["OPENAI_API_KEY"] = os.getenv("OPENAI_API_KEY")

iopex_slide_bg = os.path.join("data/Templates", "iopex_ppt_bg.png")
iopex_logo = os.path.join("data/Templates", "iopex-logo.png")
blank_bg = os.path.join("data/Templates", "blank-bg.png")

async def ask_openai(context: str, call_for: str):

    try:
        
        prompt = ""
        income_statements_prompt_template = """
            Identify all the dimensions like reporting standards, fiscal year, fiscal quarter, revenue, cost of sales, expenses, income products, service , business entity and all the metric amounts and and provide the result in the YAML file with the respective dimenstions and metrics in the following format given for one quarter and repeat for all quarters if the data is available in the context:
            Reporting Standards: 
             - GAAP
             - Non-GAAP
               - Fiscal Year: 
                - Fiscal Quarter:
                   - Quarter 1 (August)
                       - Revenue: 
                         - Product
                            - Metrics Amount
                         - Service
                            - Metrics Amount
                       - Cost of Sales
                         - Product
                            - Metrics Amount
                         - Service  
                            - Metrics Amount
                       - Gross Margins
                         - Metrics Amount
                       - Operating Income    
                         - Metrics Amount    
                       - Net Income    
                         - Metrics Amount 
                       - Net Income    
                         - Metrics Amount 
            """
        if(call_for == "generate_manifest"):
            income_statements_prompt_template = """
            Identify all the dimensions like reporting standards, fiscal year, fiscal quarter, revenue, cost of sales, expenses, income products, service , business entity and all the metric amounts and and provide the result in the YAML file with the respective dimenstions and metrics in the following format given for one quarter and repeat for all quarters if the data is available in the context:
            Reporting Standards: 
             - GAAP
             - Non-GAAP
               - Fiscal Year: 
                - Fiscal Quarter:
                   - Quarter 1 (August)
                       - Revenue: 
                         - Product
                            - Metrics Amount
                         - Service
                            - Metrics Amount
                       - Cost of Sales
                         - Product
                            - Metrics Amount
                         - Service  
                            - Metrics Amount
                       - Gross Margins
                         - Metrics Amount
                       - Operating Income    
                         - Metrics Amount    
                       - Net Income    
                         - Metrics Amount 
                       - Net Income    
                         - Metrics Amount 
            """
            prompt = income_statements_prompt_template + " from the context below" + f"\n\n{context}"
            #prompt = "Genenrate a manifest from the excel sheet in text format: " +f"\n\nExcel text: {context}"

        elif(call_for == "get_summary"):
            print("OpenAI: Generaing summary for excel..")
            summary_prompt_template = """
            For the Excel sheet in text given below, which provides details about historic financial statements, Please provide a 5-point summary. In your response, give each point as a separate sentence, and label them as (1), (2), (3), (4), and (5). All the financial numbers must be in $ / dollars"""
            prompt = summary_prompt_template + f"\n\nExcel text: {context}"
         
        response = api_openai(prompt)
        if(call_for == "generate_manifest"):
            yaml_pattern = r'```yaml(.*?)```'
            #match = re.search(yaml_pattern, str(completion.choices[0].message.content), re.DOTALL)
            match = re.search(yaml_pattern, str(response), re.DOTALL)

            if match:
                yaml_content = match.group(1).strip()
                return yaml_content
            else:
               # return str(completion.choices[0].message.content)
               return str(response)
            
        if(call_for == "get_summary"):
            #return str(completion.choices[0].message.content)
            return {"status" : 200, "summary" : str(response)}
    except Exception as e:
        return "Error from ask_open ai:" + str(e)
    
def excel_to_text(file_path: str, worksheet_name: str):

    try:
        print("opening worksheet...")
        print("filepath: " , file_path, " worksheet_name: " , worksheet_name)
        
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        worksheet_text_mapping = {}
        print("workbook loaded")
        #get excel text from the provided worksheet
        if(worksheet_name is not None):
            print("worksheet provided")
            try:
              worksheet = workbook.get_sheet_by_name(worksheet_name)
            except Exception as e:
                print(str(e))
                worksheet = workbook[worksheet_name]
            print("loaded worksheet..")
            text_content = ""
            for row in worksheet.iter_rows(values_only=True):
                row_text = " ".join(str(cell) if cell is not None else "" for cell in row)
                text_content += row_text + "\n"
            worksheet_text_mapping[worksheet_name] = text_content
        else:

        #get excel text from all worksheets
            for sheet_name in workbook.sheetnames:
                print("worksheet provided loading all worksheets..")
                worksheet = workbook[sheet_name]
                text_content = ""

                for row in worksheet.iter_rows(values_only=True):
                    row_text = " ".join(str(cell) if cell is not None else "" for cell in row)
                    text_content += row_text + "\n"

                worksheet_text_mapping[sheet_name] = text_content

        return worksheet_text_mapping
    except Exception as e:
        
        return "Error from excel_to_text : " +str(e)
    
async def generate_summary(filePath: str, sheet_name: str):

    try:
      

      worksheet_text_mapping = excel_to_text(filePath, sheet_name)
      print(worksheet_text_mapping[sheet_name])
    except Exception as e:
        err = "Error while converting excel file to text: " + str(e)
        return {"response" : err, "summary": "", "status": 400}
    
    try:

      print("calling OpenAI API..")
      response = await ask_openai(worksheet_text_mapping[sheet_name], "get_summary")
      if response["status"] == 200:
        return {"response" : "success", "summary": response["summary"], "status": 200}
    
    except Exception as e:
        err = "Error while calling OpenAI API to generate summary: " + str(e)
        return {"response" : err , "summary": "", "status": 400}

def add_footer(prs, slide, footer_text):

    slide_height = prs.slide_height
    footer = slide.shapes.add_textbox(Inches(0.25), slide_height -  Inches(1), Inches(5), Inches(0.5))
    text_frame_footer = footer.text_frame
    q = text_frame_footer.add_paragraph()
    q.text = footer_text
    q.font.size = Pt(12)
    q.font.color.rgb = RGBColor(128, 128, 128)
    q.alignment = PP_ALIGN.LEFT
 
def add_title_slide_finance(prs, Title_Text):
    
     #title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    #add backgroung picture
    left = top = Inches(0)
    width = Inches(10)
    height = Inches(7.5)
    pic = slide.shapes.add_picture(iopex_slide_bg, left, top, width, height)


    
    slide_width = prs.slide_width
    slide_height = prs.slide_height


    #header text
    left = width = height = Inches(1)
    top = Inches(3)
    title = slide.shapes.add_textbox(0, top, slide_width, height)
    title.text_frame.text = Title_Text
    text_frame_title = title.text_frame
    paragraph = text_frame_title.paragraphs[0]
    font = paragraph.font
    font.name = 'Arial'
    font.size = Pt(35)  #
    font.color.rgb = RGBColor(255, 255, 255)
    font.underline = True
    paragraph.alignment = PP_ALIGN.LEFT

    

    add_footer(prs, slide, "© All rights reserved. Confidential data")

def add_summary(prs, summary, title):
    
    
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    

    slide_width = prs.slide_width
    background_color = RGBColor(255, 255, 255)
    prs.slides[0].background.fill.solid()
    prs.slides[0].background.fill.fore_color.rgb = background_color

    # Add a background picture
    left = top = Inches(0)
    width = Inches(10)
    height = Inches(7.5)
    pic = slide.shapes.add_picture(blank_bg, left, top, width, height)

    #logo
    image_width = Inches(1)  
    image_height = Inches(0.5)  
    imgleft = slide_width - image_width  
    imgtop = Inches(0.05)
    pic = slide.shapes.add_picture(iopex_logo, imgleft, imgtop, width=image_width, height=image_height)


    # Title for the slide
    title_text = title + " Summary"
    title_left = Inches(0)  
    title_top = Inches(0)
    title_width = Inches(10)
    title_height = Inches(0.5)
    text_box = slide.shapes.add_textbox(left=title_left, top=title_top, width=title_width, height=title_height)
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    font = p.runs[0].font
    font.color.rgb = RGBColor(87, 87, 87)
    font.size = Pt(24)
    font.bold = True


    left = Inches(0)
    top = Inches(1)  
    width = Inches(5)
    height = Inches(6)  
    #text_box_left = slide.shapes.add_textbox(left, top, width, height)
    #text_frame_left = text_box_left.text_frame
    summary_content = summary
    sentences = summary_content.split("\n")
    sentences = [sentence.strip() for sentence in sentences if sentence.strip()]
    #print("sentences: ", sentences)
    
    content = slide.shapes.add_textbox(Inches(1), Inches(2), slide_width - Inches(2), Inches(5))
    
    tf = content.text_frame  
    tf.word_wrap = True
    for point in sentences:
        p = tf.add_paragraph()
        p.text = point
        p.space_after = Inches(0.25)
        p.font.size = Pt(14)
        run = p.runs[0]
        run.font.color.rgb = RGBColor(0,0,0)
    
    #add_footer(prs, slide, "© All rights reserved. Confidential data")

def get_closest_column(df, col_name):

    existing_columns = df.columns.tolist()
    
    closest_match = get_close_matches(col_name, existing_columns, n=1, cutoff=0.6)

    if closest_match:
        return df[closest_match[0]].to_numpy() 
    else:
        return None


def add_contents (prs, title, data):
    
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    background_color = RGBColor(255, 255, 255)
    prs.slides[0].background.fill.solid()
    prs.slides[0].background.fill.fore_color.rgb = background_color

    # Add a background picture
    left = top = Inches(0)
    width = Inches(10)
    height = Inches(7.5)
    pic = slide.shapes.add_picture(blank_bg, left, top, width, height)

    #logo
    image_width = Inches(1)  
    image_height = Inches(0.5)  
    imgleft = slide_width - image_width  
    imgtop = Inches(0.05)
    pic = slide.shapes.add_picture(iopex_logo, imgleft, imgtop, width=image_width, height=image_height)

    # Title for the slide
    title_text = title + " Metrics"
    title_left = Inches(0)  
    title_top = Inches(0)
    title_width = Inches(10)
    title_height = Inches(0.5)
    text_box = slide.shapes.add_textbox(left=title_left, top=title_top, width=title_width, height=title_height)
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    font = p.runs[0].font
    font.color.rgb = RGBColor(87, 87, 87)
    font.size = Pt(24)
    font.bold = True

    sub_text = "Amounts in USD"
    left = Inches(0)
    top = Inches(0.35)  
    width = Inches(10)
    height = Inches(0.5) 
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = sub_text
    p.alignment = PP_ALIGN.CENTER
    font = p.runs[0].font
    font.name = 'Arial'
    font.color.rgb = RGBColor(87, 87, 87)
    font.size = Pt(12)

    mleft = Inches(1)
    mtop = Inches(1)
    mwidth = Inches(4.5)
    mheight = Inches(5.5)
    try:
      i = 0
      
      for key, value in data.items():
        
        txBox = slide.shapes.add_textbox(mleft, mtop, mwidth, mheight)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        dimension_text = str(key)
        value_text = str(value)

        # Add the dimension text
        q_run = p.add_run()
        q_run.text = dimension_text +" : "
        q_run.font.bold = True
        q_run.font.color.rgb = RGBColor(244, 111, 3)
        
        # Add the dimension value
        month_run = p.add_run()
        month_run.text = "$" + value_text
        month_run.font.color.rgb = RGBColor(0,0,0)

        tf.word_wrap = True

        mtop += Inches(1)

    except Exception as e:
        print("error: %s" % e)


def add_meterics_data(csv_file_path, config_file, prs):
    
    try:
        print("getting dataframe from csv file..")
        df = pd.read_csv(csv_file_path)
        pd.set_option('display.float_format', lambda x: '%.4f' % x)
    except Exception as e:
        err = "Error while creating dataframe from csv: " +str(e)
        return err
    
    try:
        print("getting metrics and dimensions from config file")
        json_data = load_json(config_file)
        metrics, dimensions = extract_metrics_and_dimensions(json_data)
        for metric in metrics:
            title = metric
            data = {}
            for dimension in dimensions:
                col_header = re.sub(r'[^a-zA-Z0-9]', '', str(metric)) + "_" + re.sub(r'[^a-zA-Z0-9]', '', str(dimension))
                col_data = get_closest_column(df, col_header)
                data[dimension] = col_data[0]
            #print("title: ", title)
            add_contents(prs, title, data)

    except Exception as e:
        err = "Error while generating column header from metric and dimension: " + str(e)
        return err
async def generate_presentation(type, excel_file_path, sheet_name):
    
    #initialize varaibles
    summary = ""
    title = ""
    prs = Presentation()

    if type == "financial":
        title = "Company Financial Overview"
        try:
            
          add_title_slide_finance(prs, title)
        
        except Exception as e:
            err = "Error while adding title slide to presentation: " + str(e)
            print(err)
   

    #generate summary from excel data:
    try:
        print("generating summary from excel...")
        response = await generate_summary(excel_file_path, sheet_name)
        if response["status"] == 200:
            
            print("summary generated successfully")
            summary = response["summary"]
            #add summary slide
            try:
            #summary slide
              add_summary(prs, summary, title)
            except Exception as e:
              print("Error while adding summary slide: " + str(e))

        elif response["status"] == 400:
            return response["response"]
    
    except Exception as e:
        err = "Error while generating summary from excel: " + str(e)
        print(err)
      
    #add metric content
    try:
        excel_file_folder = excel_file_path.split("/")[-1].split(".")[0]
        path = "data/Manifest/" + excel_file_folder +"/" + sheet_name+".json"
        if(os.path.exists(path)):
            config_file_path = os.path.join("data/Manifest", excel_file_folder, sheet_name+".json")
            csv_file_path = os.path.join("data/Output", excel_file_folder, sheet_name+".csv")
            add_meterics_data(csv_file_path, config_file_path, prs)
    except Exception as e:
        err = "Error while adding metric content to presentation: " + str(e)
        print(err)

    try:

      folderName = excel_file_path.split("/")[-1].split(".")[0]
      os.makedirs(os.path.join("data/PowerPoints", folderName), exist_ok=True)
            
      #saving presentation
      ppt_name = str(sheet_name).replace(" ", "") + "_presentation.pptx"
      ppt_path = os.path.join("data/PowerPoints", folderName, ppt_name)
      prs.save(ppt_path)
      return {"ppt_path" : ppt_path,"summary" : summary, "status" : 200}
    
    except Exception as e:
      err = "Error while saving presentation: " + str(e)



