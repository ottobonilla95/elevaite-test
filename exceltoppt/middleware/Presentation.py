from functions import api_openai
from dotenv import load_dotenv

from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

import re
import os
import openai
import openpyxl
import asyncio


load_dotenv()
openai.api_key = os.getenv("OPENAI_API_KEY")
os.environ["OPENAI_API_KEY"] = os.getenv("OPENAI_API_KEY")

iopex_slide_bg = os.path.join("data/Templates", "iopex_ppt_bg.png")
iopex_logo = os.path.join("data/Templates", "iopex-logo.png")
blank_bg = os.path.join("data/Templates", "blank-bg.png")

async def ask_openai(context: str, call_for: str):

    try:
        
        prompt = ""
        income_statements_prompt_template = """
            Identify all the dimensions like reporting standards, fiscal year, fiscal quarter, revenue, cost of sales, expenses, income products, service , business entity and all the metric amounts and and provide the result in the YAML file with the respective dimenstions and metrics in the following format given for one quarter and repeat for all quarters if the data is available in the context:
            Reporting Standards: 
             - GAAP
             - Non-GAAP
               - Fiscal Year: 
                - Fiscal Quarter:
                   - Quarter 1 (August)
                       - Revenue: 
                         - Product
                            - Metrics Amount
                         - Service
                            - Metrics Amount
                       - Cost of Sales
                         - Product
                            - Metrics Amount
                         - Service  
                            - Metrics Amount
                       - Gross Margins
                         - Metrics Amount
                       - Operating Income    
                         - Metrics Amount    
                       - Net Income    
                         - Metrics Amount 
                       - Net Income    
                         - Metrics Amount 
            """
        if(call_for == "generate_manifest"):
            income_statements_prompt_template = """
            Identify all the dimensions like reporting standards, fiscal year, fiscal quarter, revenue, cost of sales, expenses, income products, service , business entity and all the metric amounts and and provide the result in the YAML file with the respective dimenstions and metrics in the following format given for one quarter and repeat for all quarters if the data is available in the context:
            Reporting Standards: 
             - GAAP
             - Non-GAAP
               - Fiscal Year: 
                - Fiscal Quarter:
                   - Quarter 1 (August)
                       - Revenue: 
                         - Product
                            - Metrics Amount
                         - Service
                            - Metrics Amount
                       - Cost of Sales
                         - Product
                            - Metrics Amount
                         - Service  
                            - Metrics Amount
                       - Gross Margins
                         - Metrics Amount
                       - Operating Income    
                         - Metrics Amount    
                       - Net Income    
                         - Metrics Amount 
                       - Net Income    
                         - Metrics Amount 
            """
            prompt = income_statements_prompt_template + " from the context below" + f"\n\n{context}"
            #prompt = "Genenrate a manifest from the excel sheet in text format: " +f"\n\nExcel text: {context}"

        elif(call_for == "get_summary"):
            print("OpenAI: Generaing summary for excel..")
            summary_prompt_template = """
            For the Excel sheet in text given below, which provides details about historic financial statements, Please provide a 5-point summary. In your response, give each point as a separate sentence, and label them as (1), (2), (3), (4), and (5). All the financial numbers must be in $ / dollars"""
            prompt = summary_prompt_template + f"\n\nExcel text: {context}"
         
        response = api_openai(prompt)
        if(call_for == "generate_manifest"):
            yaml_pattern = r'```yaml(.*?)```'
            #match = re.search(yaml_pattern, str(completion.choices[0].message.content), re.DOTALL)
            match = re.search(yaml_pattern, str(response), re.DOTALL)

            if match:
                yaml_content = match.group(1).strip()
                return yaml_content
            else:
               # return str(completion.choices[0].message.content)
               return str(response)
            
        if(call_for == "get_summary"):
            #return str(completion.choices[0].message.content)
            return {"status" : 200, "summary" : str(response)}
    except Exception as e:
        return str(e)
    
def excel_to_text(file_path: str, worksheet_name: str):

    try:
        print("opening worksheet...")
        
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        worksheet_text_mapping = {}

        #get excel text from the provided worksheet
        if(worksheet_name is not None):
            worksheet = workbook[worksheet_name]
            text_content = ""
            for row in worksheet.iter_rows(values_only=True):
                row_text = " ".join(str(cell) if cell is not None else "" for cell in row)
                text_content += row_text + "\n"
            worksheet_text_mapping[worksheet_name] = text_content
        else:

        #get excel text from all worksheets
            for sheet_name in workbook.sheetnames:
            
                worksheet = workbook[sheet_name]
                text_content = ""

                for row in worksheet.iter_rows(values_only=True):
                    row_text = " ".join(str(cell) if cell is not None else "" for cell in row)
                    text_content += row_text + "\n"

                worksheet_text_mapping[sheet_name] = text_content

        return worksheet_text_mapping
    except Exception as e:
        
        return str(e)
    
async def generate_summary(filePath: str, sheet_name: str):

    try:
      

      worksheet_text_mapping = excel_to_text(filePath, sheet_name)
    except Exception as e:
        err = "Error while converting excel file to text: " + str(e)
        return {"response" : err, "summary": "", "status": 400}
    
    try:

      print("calling OpenAI API..")
      response = await ask_openai(worksheet_text_mapping[sheet_name], "get_summary")
      if response["status"] == 200:
        return {"response" : "success", "summary": response["summary"], "status": 200}
    
    except Exception as e:
        err = "Error while calling OpenAI API to generate summary: " + str(e)
        return {"response" : err , "summary": "", "status": 400}

def add_footer(prs, slide, footer_text):

    slide_height = prs.slide_height
    footer = slide.shapes.add_textbox(Inches(0.25), slide_height -  Inches(1), Inches(5), Inches(0.5))
    text_frame_footer = footer.text_frame
    q = text_frame_footer.add_paragraph()
    q.text = footer_text
    q.font.size = Pt(12)
    q.font.color.rgb = RGBColor(128, 128, 128)
    q.alignment = PP_ALIGN.LEFT
 
def add_title_slide_finance(prs, Title_Text):
    
     #title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    #add backgroung picture
    left = top = Inches(0)
    width = Inches(10)
    height = Inches(7.5)
    pic = slide.shapes.add_picture(iopex_slide_bg, left, top, width, height)


    
    slide_width = prs.slide_width
    slide_height = prs.slide_height


    #header text
    left = width = height = Inches(1)
    top = Inches(3)
    title = slide.shapes.add_textbox(0, top, slide_width, height)
    title.text_frame.text = Title_Text
    text_frame_title = title.text_frame
    paragraph = text_frame_title.paragraphs[0]
    font = paragraph.font
    font.name = 'Arial'
    font.size = Pt(35)  #
    font.color.rgb = RGBColor(255, 255, 255)
    font.underline = True
    paragraph.alignment = PP_ALIGN.LEFT

    

    add_footer(prs, slide, "© All rights reserved. Confidential data")

def add_summary(prs, summary, title):
    
    
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    

    slide_width = prs.slide_width
    background_color = RGBColor(255, 255, 255)
    prs.slides[0].background.fill.solid()
    prs.slides[0].background.fill.fore_color.rgb = background_color

    # Add a background picture
    left = top = Inches(0)
    width = Inches(10)
    height = Inches(7.5)
    pic = slide.shapes.add_picture(blank_bg, left, top, width, height)

    #logo
    image_width = Inches(1)  
    image_height = Inches(0.5)  
    imgleft = slide_width - image_width  
    imgtop = Inches(0.05)
    pic = slide.shapes.add_picture(iopex_logo, imgleft, imgtop, width=image_width, height=image_height)


    # Title for the slide
    title_text = title + " Summary"
    title_left = Inches(0)  
    title_top = Inches(0)
    title_width = Inches(10)
    title_height = Inches(0.5)
    text_box = slide.shapes.add_textbox(left=title_left, top=title_top, width=title_width, height=title_height)
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    font = p.runs[0].font
    font.color.rgb = RGBColor(87, 87, 87)
    font.size = Pt(24)
    font.bold = True


    left = Inches(0)
    top = Inches(1)  
    width = Inches(5)
    height = Inches(6)  
    #text_box_left = slide.shapes.add_textbox(left, top, width, height)
    #text_frame_left = text_box_left.text_frame
    summary_content = summary
    sentences = summary_content.split("\n")
    sentences = [sentence.strip() for sentence in sentences if sentence.strip()]
    #print("sentences: ", sentences)
    
    content = slide.shapes.add_textbox(Inches(1), Inches(2), slide_width - Inches(2), Inches(5))
    
    tf = content.text_frame  
    tf.word_wrap = True
    for point in sentences:
        p = tf.add_paragraph()
        p.text = point
        p.space_after = Inches(0.25)
        p.font.size = Pt(14)
        run = p.runs[0]
        run.font.color.rgb = RGBColor(0,0,0)
    
    #add_footer(prs, slide, "© All rights reserved. Confidential data")

    
async def generate_presentation(type, excel_file_path, sheet_name):
    
    #initialize varaibles
    summary = ""
    title = ""
    prs = Presentation()

    if type == "financial":
        title = "Company Financial Overview"
        try:
            
          add_title_slide_finance(prs, title)
        
        except Exception as e:
            err = "Error while adding title slide to presentation: " + str(e)
            print(err)
   

    #generate summary from excel data:
    try:
        print("generating summary from excel...")
        response = await generate_summary(excel_file_path, sheet_name)
        if response["status"] == 200:
            
            print("summary generated successfully")
            summary = response["summary"]
            #add summary slide
            try:
            #summary slide
              add_summary(prs, summary, title)
            except Exception as e:
              print("Error while adding summary slide: " + str(e))

        elif response["status"] == 400:
            return response["response"]
    
    except Exception as e:
        err = "Error while generating summary from excel: " + str(e)
        print(err)

    try:

      folderName = excel_file_path.split("/")[-1].split(".")[0]
      os.makedirs(os.path.join("data/PowerPoints", folderName), exist_ok=True)
            
      #saving presentation
      ppt_name = str(sheet_name).replace(" ", "") + "_presentation.pptx"
      ppt_path = os.path.join("data/PowerPoints", folderName, ppt_name)
      prs.save(ppt_path)
      return {"ppt_path" : ppt_path,"summary" : summary, "status" : 200}
    
    except Exception as e:
      err = "Error while saving presentation: " + str(e)


