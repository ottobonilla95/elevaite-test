import openai
import re
import os
import openpyxl
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import pandas as pd
import yaml
import streamlit as st
from pptx.dml.color import RGBColor
import numpy as np

from langchain.document_loaders import CSVLoader
from langchain.indexes import VectorstoreIndexCreator
from langchain.chains import RetrievalQA
from langchain.llms import OpenAI
from langchain.embeddings.openai import OpenAIEmbeddings

current_woring_dir = os.path.dirname(os.path.realpath(__file__))
openai.api_key = ""
os.environ["OPENAI_API_KEY"] = ""
manifest_file_path = None
metrics = None
fiscal_year = None
fiscal_months = None
reporting_standards = None
csv_file_path = os.path.join(current_woring_dir, "data/Output", "output.csv")

cisco_logo = os.path.join("data/Templates", "cisco_logo.png")
cisco_slide_bg = os.path.join("data/Templates", "cisco_slide_bg.png")
cisco_toc_bg = os.path.join("data/Templates", "cisco_toc.png")
cisco_bg_metrics = os.path.join("data/Templates", "cisco_bg_metrics.png")


def ask_questions(query, chain):

    response = chain({"question": query})
    return response['result']

async def ask_openai(context: str, call_for: str):

    prompt = ""

    if(call_for == "generate_manifest"):
        income_statements_prompt_template = """
        From the Excel file in text format, identify the fiscal year, fiscal month, reporting standards, all the metrics and sub metrics and provide only YAML file in the following format:

        Fiscal Year: 
        Fiscal Months:
            - [dates]
            - [dates]
        Reporting Standards: 
        - GAAP
        ....
        Metrics: 
        - [metrics]
            - [sub metrics]
        .....
        """
        #prompt = income_statements_prompt_template + f"\n\nExcel text: {context}"
        prompt = "Genenrate a manifest from the excel sheet in text format: " +f"\n\nExcel text: {context}"

    elif(call_for == "get_summary"):
        summary_prompt_template = """
        For the Excel sheet in text given below, which provides details about historic financial statements, Please provide a 5-point summary. In your response, give each point as a separate sentence, and label them as (1), (2), (3), (4), and (5)"""
        prompt = summary_prompt_template + f"\n\nExcel text: {context}"
    
    completion = openai.ChatCompletion.create(
            model = "gpt-4",
            messages = [
                {"role":"system", "content":"You are a helpful assistant."},
                {"role":"user","content":prompt}
            ]
        )
    
    if(call_for == "generate_manifest"):
        yaml_pattern = r'```yaml(.*?)```'
        match = re.search(yaml_pattern, str(completion.choices[0].message.content), re.DOTALL)

        if match:
            yaml_content = match.group(1).strip()
            return yaml_content
        else:
            return str(completion.choices[0].message.content)
        
    if(call_for == "get_summary"):
        return str(completion.choices[0].message.content)

def excel_to_text(file_path: str, worksheet_name: str):
    workbook = openpyxl.load_workbook(file_path, data_only=True)
    worksheet_text_mapping = {}

    #get excel text from the provided worksheet
    if(worksheet_name is not None):
        worksheet = workbook[worksheet_name]
        text_content = ""
        for row in worksheet.iter_rows(values_only=True):
            row_text = " ".join(str(cell) if cell is not None else "" for cell in row)
            text_content += row_text + "\n"
        worksheet_text_mapping[worksheet_name] = text_content
    else:

    #get excel text from all worksheets
        for sheet_name in workbook.sheetnames:
        
            worksheet = workbook[sheet_name]
            text_content = ""

            for row in worksheet.iter_rows(values_only=True):
                row_text = " ".join(str(cell) if cell is not None else "" for cell in row)
                text_content += row_text + "\n"

            worksheet_text_mapping[sheet_name] = text_content

    return worksheet_text_mapping

def listFilesInDirectory(directory_path: str):
    try:
        files_list = []
        for root, dirs, files in os.walk(directory_path):
            for file in files:
                file_name = os.path.basename(file)
                files_list.append(file_name)
        return files_list
    except Exception as e:
        return str(e)
    

async def generate_manifest(file_name: str, file_path: str, save_dir: str):

    directory_path = os.path.join(current_woring_dir, save_dir, file_name)
    if os.path.exists(directory_path):
        sheet_names = listFilesInDirectory(directory_path)
        return {"message": "Manifest generated successfully!!","fileName": file_name, "sheet_names": sheet_names, "status": 200}
    
    else:
    
        os.makedirs(os.path.join(current_woring_dir, save_dir, file_name), exist_ok=True)

        worksheets_text_mapping = excel_to_text(file_path, None)
        
        sheet_names = []
        for sheet_name, text_content in worksheets_text_mapping.items():
            
            
            response = await ask_openai(text_content, "generate_manifest")

            manifest_dir = os.path.join(current_woring_dir, save_dir, file_name, f"{sheet_name}.yaml")
            sheet_name = f"{sheet_name}.yaml"
            sheet_names.append(sheet_name)

            with open(manifest_dir, 'w') as manifest_file:
                manifest_file.write(response)
        
        return {"message": "Manifest generated successfully","fileName": file_name, "sheet_names": sheet_names, "status": 200}

def getcolumns(metric):
    return list(metric.values())[0]

def get_metric_row(df, keyword):
    row, column = None, None
    for i, col in enumerate(df.columns):
        for j, value in enumerate(df[col]):
            if str(value).strip().lower() == str(keyword).lower():
                row, column = j, i
                break
    return row

def get_common_row_and_column_2(df, keyword1, keyword2):
    common_row = None
    common_col = None

    for row in df.index:
        cleaned_row_values = [str(value).strip() for value in df.loc[row].values]
        if keyword1.strip() in cleaned_row_values and keyword2.strip() in cleaned_row_values:
            common_row = row
            break

    for col in df.columns:
        cleaned_col_values = [str(value).strip() for value in df[col].values]
        #print("cleaned_col_values: ", cleaned_col_values)
        if keyword1.strip() in cleaned_col_values and keyword2.strip() in cleaned_col_values:
            common_col = col
            break
    return common_row, common_col

def Excel_to_dataframe(excel_file_path, manifest_file_path, selected_sheet):
    
    
    #load excel file into df
    df = pd.read_excel(excel_file_path, sheet_name=selected_sheet, header=None)
    df = df.fillna("")
    df = df.dropna(how='all')
    df = df.replace('\n', '', regex=True)

    #Extract YAML details
    with open(manifest_file_path, 'r') as yaml_file:
        data = yaml.safe_load(yaml_file)
    metrics = data['Metrics']
    fiscal_year = data['Fiscal Year']
    fiscal_months = data['Fiscal Months']
    reporting_standards = data['Reporting Standards']

    csv_df = pd.DataFrame()

    #constrcut csv from excel data
    for rs in reporting_standards:
        for month in fiscal_months:
            column = rs + " -[" + month.strip() + "]"
            column_values = []
            row_index = []
            row, col = get_common_row_and_column_2(df, rs, month)
            for m in metrics:
                if isinstance(m, str):
                    metric = m
                    m_row = get_metric_row(df, metric)
                    
                    if m_row is not None:
                        value = df.iat[m_row, col]
                        column_values.append(value)
                        row_index.append(metric.split('(')[0] if '(' in metric else metric)
                        #row_index.append(metric)
                if isinstance(m, dict):
                    for key, value in m.items():
                        metric = key
                        m_row = get_metric_row(df, metric)
                        
                        if m_row is not None:
                            value = df.iat[m_row, col]
                            column_values.append(value)
                            row_index.append(metric.split('(')[0] if '(' in metric else metric)
                            #row_index.append(metric)
                        sub_metric = getcolumns(m)
                        for sm in sub_metric:
                            sm_row = get_metric_row(df, sm)
                            val = df.iat[sm_row, col]
                            column_values.append(val)
                            row_index.append((metric.split('(')[0] if '(' in metric else metric) + " - " + (sm.split('(')[0] if '(' in sm else sm))
                            #row_index.append(metric + " - " + sm)
            csv_df[column] = column_values
    csv_df.index = row_index
   # st.write(csv_df)
    return csv_df

def get_df_metric(df, rowindex, colindex):
    row = df[df.index.str.contains(rowindex)]
    #print("row: ", row)
    if not row.empty and colindex in df.columns:
        # Use .loc to access the specific cell at the intersection of row and col
        value = row.loc[row.index[0], colindex]
        return value

def formatNumber(number):
    
    result = float(number) * 100
    formatted_result = "{:.1f}%".format(result)
   
    return formatted_result

def add_footer(prs, slide, footer_text):

    slide_height = prs.slide_height
    footer = slide.shapes.add_textbox(Inches(0.25), slide_height -  Inches(1), Inches(5), Inches(0.5))
    text_frame_footer = footer.text_frame
    q = text_frame_footer.add_paragraph()
    q.text = footer_text
    q.font.size = Pt(12)
    q.font.color.rgb = RGBColor(128, 128, 128)
    q.alignment = PP_ALIGN.LEFT

def add_contents_to_column (prs, text_box, contents, start_number = 1):

    tf = text_box.text_frame
    tf.word_wrap = True

    for i, content in enumerate(contents, 1):
        p = tf.add_paragraph()
        p.text = f"{start_number}. {content}"
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(20) 
        run.font.color.rgb = RGBColor(255, 255, 255)  
        run.font.italic = True
        start_number += 1

def add_toc_slide(prs, contents):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide_width = prs.slide_width
    left = top = Inches(0)
    width = Inches(10)
    height = Inches(7.5)
    pic = slide.shapes.add_picture(cisco_toc_bg, left, top, width, height)

    #logo
    image_width = Inches(1)  
    image_height = Inches(0.5)  
    imgleft = slide_width - image_width  
    imgtop = Inches(0.05)
    pic = slide.shapes.add_picture(cisco_logo, imgleft, imgtop, width=image_width, height=image_height)

    title_text = "TABLE OF CONTENTS"
    left = Inches(0)
    top = Inches(0.5)  
    width = Inches(10)
    height = Inches(0.5) 
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    font = p.runs[0].font
    font.color.rgb = RGBColor(255, 255, 255)
    font.size = Pt(24)  
    font.bold = True
   

    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4.5), Inches(5.5))

    add_contents_to_column(prs, left_box, contents[:5])
    add_contents_to_column(prs, right_box, contents[5:], start_number=len(contents[:5]) + 1)
    add_footer(prs, slide, "© Cisco and/or its affiliates. All rights reserved. Cisco Confidential")

def add_summary(prs, summary):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide_width = prs.slide_width

    # Add a background picture
    left = top = Inches(0)
    width = Inches(10)
    height = Inches(7.5)
    pic = slide.shapes.add_picture(cisco_slide_bg, left, top, width, height)

    #cisco logo
    image_width = Inches(1)  
    image_height = Inches(0.5)  
    imgleft = slide_width - image_width  
    imgtop = Inches(0.05)
    pic = slide.shapes.add_picture(cisco_logo, imgleft, imgtop, width=image_width, height=image_height)

    # Title for the slide
    title_text = "SUMMARY"
    left = Inches(0)
    top = Inches(0.5)
    width = Inches(10)
    height = Inches(0.5)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    font = p.runs[0].font
    font.color.rgb = RGBColor(255, 255, 255)
    font.size = Pt(24)
    font.bold = True


    left = Inches(0)
    top = Inches(1.25)  
    width = Inches(5)
    height = Inches(6)  
    text_box_left = slide.shapes.add_textbox(left, top, width, height)
    text_frame_left = text_box_left.text_frame
    summary_content = summary
    sentences = summary_content.split("\n")
    
    content = slide.shapes.add_textbox(Inches(1), Inches(2), slide_width - Inches(2), Inches(5))
    
    tf = content.text_frame  # Move this outside of the loop to create a single text frame for the entire content
    tf.word_wrap = True
    for point in sentences:
        p = tf.add_paragraph()
        p.text = point
        p.space_after = Inches(0.25)
        p.font.size = Pt(14)
        run = p.runs[0]
        run.font.color.rgb = RGBColor(255, 255, 255)
    
    add_footer(prs, slide, "© Cisco and/or its affiliates. All rights reserved. Cisco Confidential")

def add_title_slide_cisco_finance(prs, Title_Text, fiscal_year):
     #title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    background_color = RGBColor(13, 39, 77)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    prs.slides[0].background.fill.solid()
    prs.slides[0].background.fill.fore_color.rgb = background_color

    #cisco logo
    image_width = Inches(1)  
    image_height = Inches(0.5)  
    imgleft = slide_width - image_width  
    imgtop = Inches(0.05)
    pic = slide.shapes.add_picture(cisco_logo, imgleft, imgtop, width=image_width, height=image_height)

    #header text
    left = width = height = Inches(1)
    top = Inches(2)
    wearecisco = slide.shapes.add_textbox(0, top, slide_width, height)
    wearecisco.text_frame.text = "#WEARECISCO"
    text_frame_wearecisco = wearecisco.text_frame
    paragraph = text_frame_wearecisco.paragraphs[0]
    font = paragraph.font
    font.name = 'Arial'
    font.size = Pt(35)  #
    font.color.rgb = RGBColor(0, 187, 235)
    font.underline = True
    paragraph.alignment = PP_ALIGN.LEFT

    #header text
    title = slide.shapes.add_textbox(0, Inches(3), slide_width, height)
    title.text_frame.text = Title_Text
    text_frame_title = title.text_frame
    paragraph = text_frame_title.paragraphs[0]
    font = paragraph.font
    font.name = 'Arial'
    font.size = Pt(30)
    font.color.rgb = RGBColor(0, 187, 235)
    paragraph.alignment = PP_ALIGN.LEFT
    sub_title = slide.shapes.add_textbox(0, Inches(4), slide_width, height)
    sub_title.text_frame.text = fiscal_year
    text_frame_sub_title = sub_title.text_frame
    paragraph = text_frame_sub_title.paragraphs[0]
    font = paragraph.font
    font.name = 'Arial'
    font.size = Pt(20)
    font.color.rgb = RGBColor(255, 255, 255)
    paragraph.alignment = PP_ALIGN.LEFT

    add_footer(prs, slide, "© Cisco and/or its affiliates. All rights reserved. Cisco Confidential")

def add_excel_contents (prs, metrics, rs, fiscal_months, sub_metrics, df):

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide_width = prs.slide_width
    left = top = Inches(0)
    width = Inches(10)
    height = Inches(7.5)
    pic = slide.shapes.add_picture(cisco_bg_metrics, left, top, width, height)

    #logo
    image_width = Inches(1)  
    image_height = Inches(0.5)  
    imgleft = slide_width - image_width  
    imgtop = Inches(0.05)
    pic = slide.shapes.add_picture(cisco_logo, imgleft, imgtop, width=image_width, height=image_height)

    title_text = metrics
    left = Inches(0)
    top = Inches(0.5)  
    width = Inches(10)
    height = Inches(0.5) 
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    font = p.runs[0].font
    font.name = 'Arial'
    font.color.rgb = RGBColor(0, 187, 235)
    font.size = Pt(24)  
    font.bold = True

    subtitle_text = rs
    left = Inches(0)
    top = Inches(0.95)  
    width = Inches(10)
    height = Inches(0.5) 
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = subtitle_text
    p.alignment = PP_ALIGN.CENTER
    font = p.runs[0].font
    font.name = 'Arial'
    font.color.rgb = RGBColor(255, 255, 255)
    font.size = Pt(20)  
   

    mleft = Inches(1)
    mtop = Inches(2)
    mwidth = Inches(4.5)
    mheight = Inches(5.5)

    i = 0
    for month in fiscal_months:
        if i < 2:
            mleft = Inches(0.5 + i * 5)
            mtop = mtop
            i += 1
        else:
            mleft = Inches(0.5 + (i - 2) * 5)
            mtop = Inches(5)
            i += 1

        txBox = slide.shapes.add_textbox(mleft, mtop, mwidth, mheight)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        q_text = f"Q{i} : "
        month_text = month

        # Add the Q text
        q_run = p.add_run()
        q_run.text = q_text
        q_run.font.bold = True
        q_run.font.color.rgb = RGBColor(0, 187, 235)

        # Add the month text
        month_run = p.add_run()
        month_run.text = month_text
        month_run.font.color.rgb = RGBColor(255, 255, 255)

        tf.word_wrap = True
         
        
        txBox = slide.shapes.add_textbox(mleft, mtop, mwidth, mheight)
        tf = txBox.text_frame
        p = tf.add_paragraph()

        # Create separate runs for different parts of the text
        q_text = f"Q{i} : "
        q_run.text = q_text
        q_run.font.bold = True
        month_text = month
        month_run.text = month_text

        # Add the entire text
        full_run = p.add_run()
        full_text = f"{q_text}{month_text}"
        full_run.text = full_text

        # Apply underline and blue color to the entire text
        full_run.font.underline = True
        full_run.font.color.rgb = RGBColor(0, 187, 235)

        

        tf.word_wrap = True
        for sm in sub_metrics:
            rowindex = metrics + " - " + (sm.split('(')[0] if '(' in sm else sm)
            colindex =  rs + " -[" + month.strip() + "]"
            val = get_df_metric(df, rowindex, colindex)
            
            try:
                val = formatNumber(val) if str(np.float64(val)).startswith("0.") else val
            except Exception as e:
                                val = ""
        
            sm_val_text = f"{sm.split('(')[0] if '(' in sm else sm} : {val}"
            p = tf.add_paragraph()
            p.text = sm_val_text
            p.font.color.rgb = RGBColor(255, 255, 255)

async def generate_summary(filePath: str, sheet_name: str):

    worksheet_text_mapping = excel_to_text(filePath, sheet_name)

    response = await ask_openai(worksheet_text_mapping[sheet_name], "get_summary")

    return {"summary": response, "status": 200}

def generate_cisco_presentation(excel_file_path, manifest_file_path, summary, selected_sheet):
    
    #Extract YAML details
    with open(manifest_file_path, 'r') as yaml_file:
        data = yaml.safe_load(yaml_file)
    metrics = data['Metrics']
    fiscal_year = data['Fiscal Year']
    fiscal_months = data['Fiscal Months']
    reporting_standards = data['Reporting Standards']

    #Convert Excel to Datafeme/csv
    df = Excel_to_dataframe(excel_file_path, manifest_file_path, selected_sheet)

    #create table of contents list
    contents = ["Summary"]

    for content in metrics:
        if isinstance(content, dict):
            for key, value in content.items():
                contents.append(key.split('(')[0] if '(' in key else key)


    contents.append("Miscellaneous")

    prs = Presentation()

    #title slide
    add_title_slide_cisco_finance(prs, "CISCO: Company Financial Overview", fiscal_year)

    #table of contents
    add_toc_slide(prs, contents)

    #summary slide
    add_summary(prs, summary)

    #excel content slides
    for m in metrics:
        for rs in reporting_standards:
            if isinstance(m, dict):
                for key, value in m.items():
                    metric = key.split('(')[0] if '(' in key else key
                    sub_metric = getcolumns(m)
                    add_excel_contents(prs, metric, rs, fiscal_months, sub_metric, df)
            

    #saving presentation
    ppt_path = os.path.join("data/PowerPoints", "cisco_presentation.pptx")
    prs.save(ppt_path)
    return ppt_path

