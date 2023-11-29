'''
Functions pretaining to
    1. Generating Manifest from Excel Workbook
    2. Summary generation from Excel Workbook data using openai
    3. Presentation generation from Excel Workbook 
'''

import os
import openpyxl
import openai
from openai import OpenAI

from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from dotenv import load_dotenv


current_woring_dir = os.path.dirname(os.path.realpath(__file__))
load_dotenv()
openai.api_key = os.getenv("OPENAI_API_KEY")
os.environ["OPENAI_API_KEY"] = os.getenv("OPENAI_API_KEY")

manifest_file_path = None
metrics = None
fiscal_year = None
fiscal_months = None
reporting_standards = None
csv_file_path = os.path.join(current_woring_dir, "data/Output", "output.csv")

cisco_logo = os.path.join("data/Templates", "cisco_logo.png")
cisco_slide_bg = os.path.join("data/Templates", "cisco_slide_bg.png")
cisco_toc_bg = os.path.join("data/Templates", "cisco_toc.png")
cisco_bg_metrics = os.path.join("data/Templates", "cisco_bg_metrics.png")

def listFilesInDirectory(directory_path: str):
    try:
        files_list = []
        for root, dirs, files in os.walk(directory_path):
            for file in files:
                file_name = os.path.basename(file)
                files_list.append(file_name)
        return files_list
    except Exception as e:
        return str(e)
    
def excel_to_text(file_path: str, worksheet_name: str):

    try:
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        worksheet_text_mapping = {}

        #get excel text from the provided worksheet
        if(worksheet_name is not None):
            worksheet = workbook[worksheet_name]
            text_content = ""
            for row in worksheet.iter_rows(values_only=True):
                row_text = " ".join(str(cell) if cell is not None else "" for cell in row)
                text_content += row_text + "\n"
            worksheet_text_mapping[worksheet_name] = text_content
        else:

        #get excel text from all worksheets
            for sheet_name in workbook.sheetnames:
            
                worksheet = workbook[sheet_name]
                text_content = ""

                for row in worksheet.iter_rows(values_only=True):
                    row_text = " ".join(str(cell) if cell is not None else "" for cell in row)
                    text_content += row_text + "\n"

                worksheet_text_mapping[sheet_name] = text_content

        return worksheet_text_mapping
    except Exception as e:
        
        return str(e)

async def ask_openai(context: str, call_for: str):

    try:
        prompt = ""
        income_statements_prompt_template = """
            Identify all the dimensions like reporting standards, fiscal year, fiscal quarter, revenue, cost of sales, expenses, income products, service , business entity and all the metric amounts and and provide the result in the YAML file with the respective dimenstions and metrics in the following format given for one quarter and repeat for all quarters if the data is available in the context:
            Reporting Standards: 
             - GAAP
             - Non-GAAP
               - Fiscal Year: 
                - Fiscal Quarter:
                   - Quarter 1 (August)
                       - Revenue: 
                         - Product
                            - Metrics Amount
                         - Service
                            - Metrics Amount
                       - Cost of Sales
                         - Product
                            - Metrics Amount
                         - Service  
                            - Metrics Amount
                       - Gross Margins
                         - Metrics Amount
                       - Operating Income    
                         - Metrics Amount    
                       - Net Income    
                         - Metrics Amount 
                       - Net Income    
                         - Metrics Amount 
            """
        if(call_for == "generate_manifest"):
            income_statements_prompt_template = """
            Identify all the dimensions like reporting standards, fiscal year, fiscal quarter, revenue, cost of sales, expenses, income products, service , business entity and all the metric amounts and and provide the result in the YAML file with the respective dimenstions and metrics in the following format given for one quarter and repeat for all quarters if the data is available in the context:
            Reporting Standards: 
             - GAAP
             - Non-GAAP
               - Fiscal Year: 
                - Fiscal Quarter:
                   - Quarter 1 (August)
                       - Revenue: 
                         - Product
                            - Metrics Amount
                         - Service
                            - Metrics Amount
                       - Cost of Sales
                         - Product
                            - Metrics Amount
                         - Service  
                            - Metrics Amount
                       - Gross Margins
                         - Metrics Amount
                       - Operating Income    
                         - Metrics Amount    
                       - Net Income    
                         - Metrics Amount 
                       - Net Income    
                         - Metrics Amount 
            """
            prompt = income_statements_prompt_template + " from the context below" + f"\n\n{context}"
            #prompt = "Genenrate a manifest from the excel sheet in text format: " +f"\n\nExcel text: {context}"

        elif(call_for == "get_summary"):
            summary_prompt_template = """
            For the Excel sheet in text given below, which provides details about historic financial statements, Please provide a 5-point summary. In your response, give each point as a separate sentence, and label them as (1), (2), (3), (4), and (5). All the financial numbers must be in $ / dollars"""
            prompt = summary_prompt_template + f"\n\nExcel text: {context}"
         
        response = api_openai(prompt)
        if(call_for == "generate_manifest"):
            yaml_pattern = r'```yaml(.*?)```'
            #match = re.search(yaml_pattern, str(completion.choices[0].message.content), re.DOTALL)
            match = re.search(yaml_pattern, str(response), re.DOTALL)

            if match:
                yaml_content = match.group(1).strip()
                return yaml_content
            else:
               # return str(completion.choices[0].message.content)
               return str(response)
            
        if(call_for == "get_summary"):
            #return str(completion.choices[0].message.content)
            return str(response)
    except Exception as e:
        return str(e)

def api_openai(content: str):
    client = OpenAI()
    chat_completion = client.chat.completions.create(
        messages=[
            {"role":"system", "content":"You are a helpful assistant."},
            {"role":"user","content":content}

        ],
        model="gpt-4",
    )
    return(chat_completion.choices[0].message.content)

async def generate_manifest(file_name: str, file_path: str, save_dir: str):

    directory_path = os.path.join(current_woring_dir, save_dir, file_name)
    if os.path.exists(directory_path):
        sheet_names = listFilesInDirectory(directory_path)
        return {"message": "Manifest generated successfully!!","fileName": file_name, "sheet_names": sheet_names, "status": 200}
    
    else:
    
        os.makedirs(os.path.join(current_woring_dir, save_dir, file_name), exist_ok=True)

        worksheets_text_mapping = excel_to_text(file_path, None)
        
        sheet_names = []
        for sheet_name, text_content in worksheets_text_mapping.items():
            
            
            response = await ask_openai(text_content, "generate_manifest")

            manifest_dir = os.path.join(current_woring_dir, save_dir, file_name, f"{sheet_name}.yaml")
            sheet_name = f"{sheet_name}.yaml"
            sheet_names.append(sheet_name)

            with open(manifest_dir, 'w') as manifest_file:
                manifest_file.write(response)
        
        return {"message": "Manifest generated successfully","fileName": file_name, "sheet_names": sheet_names, "status": 200}

def add_footer(prs, slide, footer_text):

    slide_height = prs.slide_height
    footer = slide.shapes.add_textbox(Inches(0.25), slide_height -  Inches(1), Inches(5), Inches(0.5))
    text_frame_footer = footer.text_frame
    q = text_frame_footer.add_paragraph()
    q.text = footer_text
    q.font.size = Pt(12)
    q.font.color.rgb = RGBColor(128, 128, 128)
    q.alignment = PP_ALIGN.LEFT

def add_title_slide_cisco_finance(prs, Title_Text, fiscal_year):
     #title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    background_color = RGBColor(13, 39, 77)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    prs.slides[0].background.fill.solid()
    prs.slides[0].background.fill.fore_color.rgb = background_color

    #cisco logo
    image_width = Inches(1)  
    image_height = Inches(0.5)  
    imgleft = slide_width - image_width  
    imgtop = Inches(0.05)
    pic = slide.shapes.add_picture(cisco_logo, imgleft, imgtop, width=image_width, height=image_height)

    #header text
    left = width = height = Inches(1)
    top = Inches(2)
    wearecisco = slide.shapes.add_textbox(0, top, slide_width, height)
    wearecisco.text_frame.text = "#WEARECISCO"
    text_frame_wearecisco = wearecisco.text_frame
    paragraph = text_frame_wearecisco.paragraphs[0]
    font = paragraph.font
    font.name = 'Arial'
    font.size = Pt(35)  #
    font.color.rgb = RGBColor(0, 187, 235)
    font.underline = True
    paragraph.alignment = PP_ALIGN.LEFT

    #header text
    title = slide.shapes.add_textbox(0, Inches(3), slide_width, height)
    title.text_frame.text = Title_Text
    text_frame_title = title.text_frame
    paragraph = text_frame_title.paragraphs[0]
    font = paragraph.font
    font.name = 'Arial'
    font.size = Pt(30)
    font.color.rgb = RGBColor(0, 187, 235)
    paragraph.alignment = PP_ALIGN.LEFT
    sub_title = slide.shapes.add_textbox(0, Inches(4), slide_width, height)
    sub_title.text_frame.text = fiscal_year
    text_frame_sub_title = sub_title.text_frame
    paragraph = text_frame_sub_title.paragraphs[0]
    font = paragraph.font
    font.name = 'Arial'
    font.size = Pt(20)
    font.color.rgb = RGBColor(255, 255, 255)
    paragraph.alignment = PP_ALIGN.LEFT

    add_footer(prs, slide, "© Cisco and/or its affiliates. All rights reserved. Cisco Confidential")

def add_contents_to_column (prs, text_box, contents, start_number = 1):

    tf = text_box.text_frame
    tf.word_wrap = True

    for i, content in enumerate(contents, 1):
        p = tf.add_paragraph()
        p.text = f"{start_number}. {content}"
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(20) 
        run.font.color.rgb = RGBColor(255, 255, 255)  
        run.font.italic = True
        start_number += 1

def add_toc_slide(prs, contents):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide_width = prs.slide_width
    left = top = Inches(0)
    width = Inches(10)
    height = Inches(7.5)
    pic = slide.shapes.add_picture(cisco_toc_bg, left, top, width, height)

    #logo
    image_width = Inches(1)  
    image_height = Inches(0.5)  
    imgleft = slide_width - image_width  
    imgtop = Inches(0.05)
    pic = slide.shapes.add_picture(cisco_logo, imgleft, imgtop, width=image_width, height=image_height)

    title_text = "TABLE OF CONTENTS"
    left = Inches(0)
    top = Inches(0.5)  
    width = Inches(10)
    height = Inches(0.5) 
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    font = p.runs[0].font
    font.color.rgb = RGBColor(255, 255, 255)
    font.size = Pt(24)  
    font.bold = True
   

    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4.5), Inches(5.5))

    if contents:
        total_items = len(contents)
        mid_point = total_items // 2

        add_contents_to_column(prs, left_box, contents[:mid_point])

        if total_items % 2 == 0:
            add_contents_to_column(prs, right_box, contents[mid_point:], start_number=mid_point + 1)
        else:
            add_contents_to_column(prs, right_box, contents[mid_point + 1:], start_number=mid_point + 2)   
    '''add_contents_to_column(prs, left_box, contents[:5])
    add_contents_to_column(prs, right_box, contents[5:], start_number=len(contents[:5]) + 1)'''
    add_footer(prs, slide, "© Cisco and/or its affiliates. All rights reserved. Cisco Confidential")

def add_summary(prs, summary):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide_width = prs.slide_width

    # Add a background picture
    left = top = Inches(0)
    width = Inches(10)
    height = Inches(7.5)
    pic = slide.shapes.add_picture(cisco_slide_bg, left, top, width, height)

    #cisco logo
    image_width = Inches(1)  
    image_height = Inches(0.5)  
    imgleft = slide_width - image_width  
    imgtop = Inches(0.05)
    pic = slide.shapes.add_picture(cisco_logo, imgleft, imgtop, width=image_width, height=image_height)

    # Title for the slide
    title_text = "SUMMARY"
    left = Inches(0)
    top = Inches(0.5)
    width = Inches(10)
    height = Inches(0.5)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    font = p.runs[0].font
    font.color.rgb = RGBColor(255, 255, 255)
    font.size = Pt(24)
    font.bold = True


    left = Inches(0)
    top = Inches(1.25)  
    width = Inches(5)
    height = Inches(6)  
    text_box_left = slide.shapes.add_textbox(left, top, width, height)
    text_frame_left = text_box_left.text_frame
    summary_content = summary
    sentences = summary_content.split("\n")
    
    content = slide.shapes.add_textbox(Inches(1), Inches(2), slide_width - Inches(2), Inches(5))
    
    tf = content.text_frame  # Move this outside of the loop to create a single text frame for the entire content
    tf.word_wrap = True
    for point in sentences:
        p = tf.add_paragraph()
        p.text = point
        p.space_after = Inches(0.25)
        p.font.size = Pt(14)
        run = p.runs[0]
        run.font.color.rgb = RGBColor(255, 255, 255)
    
    add_footer(prs, slide, "© Cisco and/or its affiliates. All rights reserved. Cisco Confidential")

async def generate_summary(filePath: str, sheet_name: str):

    try:
       
        worksheet_text_mapping = excel_to_text(filePath, sheet_name)
        

        response = await ask_openai(worksheet_text_mapping[sheet_name], "get_summary")

        return {"summary": response, "status": 200}
    except Exception as e:
        return str(e)
    
def generate_presentation(excel_file_path, manifest_file_path, summary, selected_sheet):
    
    try:
        #create table of contents list
        contents = ["Summary", "Company Overview", "Financial Overview", "Miscellaneous", "Conclusion"]

        prs = Presentation()
        try:
            #title slide
            add_title_slide_cisco_finance(prs, "CISCO: Company Financial Overview", "FY 2023")
        except Exception as e:
            print("Error while adding title slide: " + str(e))

        try:
            #table of contents
            add_toc_slide(prs, contents)
        except Exception as e:
            print("Error while adding table of content slide: " + str(e))

        try:
            #summary slide
            add_summary(prs, summary)
        except Exception as e:
            print("Error while adding summary slide: " + str(e))

        try:
            folderName = excel_file_path.split("/")[-1].split(".")[0]

            os.makedirs(os.path.join("data/PowerPoints", folderName), exist_ok=True)
            #saving presentation
            ppt_name = str(selected_sheet).replace(" ", "") + "_presentation.pptx"
            ppt_path = os.path.join("data/PowerPoints", folderName, ppt_name)
            prs.save(ppt_path)
            return ppt_path
        except Exception as e:
            folderName = "Excel_Sheets"
            os.makedirs(os.path.join("data/PowerPoints", folderName), exist_ok=True)
            #saving presentation
            ppt_name = str(selected_sheet).replace(" ", "") + "_presentation.pptx"
            pdf_name = str(selected_sheet).replace(" ", "") + "_presentation.pdf"
            ppt_path = os.path.join("data/PowerPoints", folderName, ppt_name)
            prs.save(ppt_path)
            prs.save(pdf_name)
            return ppt_path
    except Exception as e:
        return("Error while generating PowerPoints: " + str(e)) 